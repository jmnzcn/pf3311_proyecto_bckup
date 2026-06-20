#!/usr/bin/env python3
"""Construye presentación PowerPoint del estudio."""

from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from informe_data_local import AnalysisBundle
from informe_narrative_local import InformeNarrative

UCR_BLUE = RGBColor(0, 56, 101)
TEXT_DARK = RGBColor(33, 37, 41)
TEXT_MUTED = RGBColor(90, 98, 104)

TITLE = "Efecto de la visibilidad de un agente virtual"
SUBTITLE = "Desempeño y confianza en tareas de decisión · PF-3311"
AUTHOR = "Ney Fred Jiménez Campos (B03230)"
AFFILIATION = "UCR · Posgrado en Computación e Informática"


def _set_title_style(shape, *, size: int = 32, color: RGBColor = UCR_BLUE) -> None:
    if not shape or not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.bold = True
            run.font.color.rgb = color


def _add_bullets(text_frame, lines: list[str]) -> None:
    if not lines:
        text_frame.text = ""
        return
    text_frame.text = lines[0]
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = TEXT_DARK
    for line in lines[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = TEXT_DARK


def _slide_title_content(prs: Presentation, title: str, bullets: list[str]):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    _set_title_style(slide.shapes.title, size=28)
    body = slide.placeholders[1]
    _add_bullets(body.text_frame, bullets)
    return slide


def _slide_figure(prs: Presentation, title: str, image_path: Path, caption: str = "") -> None:
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    _set_title_style(slide.shapes.title, size=26)
    if image_path.is_file():
        slide.shapes.add_picture(
            str(image_path),
            Inches(1.0),
            Inches(1.45),
            width=Inches(8.0),
        )
    if caption:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(6.55), Inches(8.4), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER


def _group_means_bullets(rows: list[dict[str, str]], value_key: str, suffix: str = "") -> list[str]:
    bullets: list[str] = []
    for row in rows:
        cond = row.get("Condition", "")
        val = row.get(value_key, "")
        n = row.get("NParticipants", "")
        if cond and val:
            bullets.append(f"Condición {cond}: {val}{suffix} (n={n})")
    return bullets or ["Sin datos exportados."]


def _integrity_bullets(bundle: AnalysisBundle) -> list[str]:
    integrity = bundle.integrity_map()
    if not integrity:
        return ["Sin resumen de viabilidad en _analysis/."]
    return [
        f"Integridad CSV válida: {integrity.get('ValidPct', 'n/a')}%",
        f"Latencia Gemini ≤5 s (meta 90%): {'cumple' if integrity.get('GeminiPasses90Pct') == '1' else 'no cumple'}",
        f"TTS condición C ≥85%: {'cumple' if integrity.get('TtsPasses85Pct') == '1' else 'no cumple'}",
        f"Filtraciones del modelo: {integrity.get('TotalModelLeaks', '0')}",
    ]


def _chat_bullets(bundle: AnalysisBundle) -> list[str]:
    if not bundle.chat_group:
        return ["Sin métricas agregadas de chat B vs C."]
    lines: list[str] = []
    for row in bundle.chat_group:
        metric = row.get("Metric", "")
        if not metric:
            continue
        lines.append(
            f"{metric}: B={row.get('Mean_B', '')} | C={row.get('Mean_C', '')}"
        )
    return lines or ["Sin datos de agente."]


def build_pptx(
    bundle: AnalysisBundle,
    narrative: InformeNarrative,
    output_path: Path,
) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Portada
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = TITLE
    slide.placeholders[1].text = (
        f"{SUBTITLE}\n{AUTHOR}\n{AFFILIATION}\n"
        f"{datetime.now().strftime('%d/%m/%Y')}"
    )
    _set_title_style(slide.shapes.title, size=30)
    for paragraph in slide.placeholders[1].text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = TEXT_DARK

    n_canon = len(bundle.canonical_sessions)
    n_complete = len(bundle.complete_canonical)

    _slide_title_content(
        prs,
        "Contexto y preguntas de investigación",
        [
            "Prototipo Unity · condiciones A (sin agente), B (chat), C (avatar + voz).",
            "RQ1: ¿Difieren los aciertos entre condiciones?",
            "RQ2: ¿Difieren confianza y percepción (Unity + Forms)?",
            "RQ3: ¿Cómo varía la calibración confianza–precisión?",
            "Hipótesis exploratorias: C ≈ B > A en precisión; mayor confianza en C.",
        ],
    )

    _slide_title_content(
        prs,
        "Diseño del estudio",
        [
            "Within-subjects: cada participante completa A, B y C (6 ítems por bloque).",
            f"Sesiones canónicas analizadas: n={n_canon}.",
            f"Completitud A+B+C: {n_complete}/{n_canon}.",
            "Confianza 1–7 por ítem (Unity); RAW-TLX y meCUE tras cada bloque.",
            "Agente: Gemini 2.5 Flash (B/C); Azure TTS + avatar en C.",
        ],
    )

    perfil_lines = narrative.sections.get("perfil", [])[:5]
    if perfil_lines:
        _slide_title_content(prs, "Características de la muestra (Form 0)", perfil_lines)

    _slide_figure(
        prs,
        "RQ1 — Precisión (% aciertos)",
        bundle.figures_dir / "rq1_precision.png",
        "Medias grupales por condición",
    )
    _slide_title_content(
        prs,
        "RQ1 — Resumen numérico",
        _group_means_bullets(bundle.rq1_group, "MeanAccuracyPct", "%"),
    )

    _slide_figure(
        prs,
        "RQ2 — Confianza (Unity, 1–7)",
        bundle.figures_dir / "rq2_confidence.png",
    )
    rq2_bullets = _group_means_bullets(bundle.rq2_group, "MeanConfidence")
    if bundle.rq2_forms_tlx:
        rq2_bullets.append("RAW-TLX y meCUE: ver exportaciones Forms en el informe.")
    _slide_title_content(prs, "RQ2 — Confianza y cuestionarios", rq2_bullets)

    _slide_figure(
        prs,
        "RQ3 — Brecha de calibración",
        bundle.figures_dir / "rq3_calibration_gap.png",
        "Confianza normalizada − precisión del bloque",
    )

    chat_figs = (
        ("chat_helpscore_b_vs_c.png", "Agente — HelpScore B vs C"),
        ("chat_exchanges_b_vs_c.png", "Agente — Intercambios de chat"),
    )
    for fig_name, fig_title in chat_figs:
        path = bundle.figures_dir / fig_name
        if path.is_file():
            _slide_figure(prs, fig_title, path)

    _slide_title_content(prs, "Agente conversacional (B vs C)", _chat_bullets(bundle))

    _slide_title_content(prs, "Viabilidad técnica del estudio", _integrity_bullets(bundle))

    respuestas = narrative.sections.get("respuestas_rq", [])
    if respuestas:
        _slide_title_content(
            prs,
            "Respuestas a las preguntas de investigación",
            respuestas[1:4] if len(respuestas) > 3 else respuestas,
        )

    conclusiones = narrative.sections.get("conclusiones", [])
    discusion = narrative.sections.get("discusion", [])[-2:]
    _slide_title_content(
        prs,
        "Conclusiones y limitaciones",
        conclusiones[:5] + discusion,
    )

    _slide_title_content(
        prs,
        "Cierre",
        [
            "PF-3311 — Agentes Virtuales Inteligentes",
            AUTHOR,
            "Datos y figuras: carpeta _analysis/ del pipeline principal",
            "Documentos locales: informe PDF, artículo paper, esta presentación",
        ],
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
